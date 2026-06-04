"""
generation.py – slide-spec building and rendering for the myHQ proposal tool.

Exported:
    build_london_slides(proposal, db_centres, manual_centres) -> list[list[dict]]
    build_india_slides(proposal,  db_centres, manual_centres) -> list[list[dict]]
    render_pptx(slides, out_path) -> out_path
    render_pdf(slides,  out_path) -> out_path
    get_logo_png()                -> path str

All coordinates are in EMU (914 400 EMU = 1 inch).
Canvas: W=12 192 000  H=6 858 000.
"""

import os
import json
import subprocess
import datetime

BASE_DIR   = os.path.dirname(os.path.abspath(__file__))
IMG_DIR    = os.path.join(BASE_DIR, 'static', 'images')
COVER_IMG  = os.path.join(IMG_DIR, 'cover-city.jpg')      # city-of-london high-res cover
CROP_IMG   = os.path.join(IMG_DIR, 'workspace-helps.jpg')  # portrait photo for How myHQ Helps
MAP_IMG    = os.path.join(IMG_DIR, 'london-map.png')
TUBE_DIR   = os.path.join(IMG_DIR, 'tube')

# Station name → primary tube line badge key
STATION_LINE_MAP = {
    # A
    'aldgate':                  'circle',
    'aldgate east':             'district',
    'angel':                    'northern',
    'archway':                  'northern',
    # B
    'baker street':             'jubilee',
    'bank':                     'central',
    'barbican':                 'circle',
    'bayswater':                'circle',
    'bermondsey':               'jubilee',
    'bethnal green':            'central',
    'blackfriars':              'circle',
    'bond street':              'jubilee',
    'borough':                  'northern',
    'brixton':                  'victoria',
    # C
    'caledonian road':          'piccadilly',
    'camden town':              'northern',
    'canada water':             'jubilee',
    'canary wharf':             'jubilee',
    'cannon street':            'circle',
    'chancery lane':            'central',
    'chiswick':                 'district',
    'city thameslink':          'overground',
    'clapham common':           'northern',
    'clapham junction':         'overground',
    'clapham north':            'northern',
    'clapham south':            'northern',
    'cockfosters':              'piccadilly',
    'covent garden':            'piccadilly',
    'crossharbour':             'dlr',
    # D
    'dalston junction':         'overground',
    'elephant & castle':        'northern',
    'elephant and castle':      'northern',
    'euston':                   'northern',
    'euston square':            'circle',
    # F
    'farringdon':               'elizabeth',
    'finsbury park':            'piccadilly',
    # G
    'goodge street':            'northern',
    'green park':               'jubilee',
    'greenwich':                'dlr',
    # H
    'hackney central':          'overground',
    'hammersmith':              'district',
    'heron quays':              'dlr',
    'highbury & islington':     'victoria',
    'highbury and islington':   'victoria',
    'holborn':                  'central',
    # K
    "king's cross":             'victoria',
    "king's cross st pancras":  'victoria',
    'kings cross':              'victoria',
    'knightsbridge':            'piccadilly',
    # L
    'lambeth north':            'bakerloo',
    'leicester square':         'northern',
    'liverpool street':         'elizabeth',
    'london bridge':            'jubilee',
    # M
    'mansion house':            'circle',
    'marble arch':              'central',
    'marylebone':               'bakerloo',
    'mile end':                 'central',
    'monument':                 'circle',
    'moorgate':                 'circle',
    'mornington crescent':      'northern',
    # N
    'new cross':                'overground',
    # O
    'old street':               'northern',
    'oxford circus':            'victoria',
    # P
    'paddington':               'elizabeth',
    'peckham rye':              'overground',
    'pimlico':                  'victoria',
    'piccadilly circus':        'piccadilly',
    # R
    'ravenscourt park':         'district',
    'richmond':                 'district',
    'royal oak':                'hammersmith',
    'russell square':           'piccadilly',
    # S
    'shadwell':                 'dlr',
    'shoreditch high street':   'overground',
    'sloane square':            'circle',
    'south kensington':         'piccadilly',
    'south quay':               'dlr',
    'southwark':                'jubilee',
    'st james\'s park':         'circle',
    'st james park':            'circle',
    'st paul\'s':               'central',
    'st pauls':                 'central',
    'stamford brook':           'district',
    'stepney green':            'district',
    'stockwell':                'victoria',
    'stratford':                'elizabeth',
    # T
    'temple':                   'circle',
    'tottenham court road':     'elizabeth',
    'tower hill':               'circle',
    'tower gateway':            'dlr',
    # G
    'great portland street':    'circle',
    'gloucester road':          'circle',
    'goodge street':            'northern',
    # V
    'vauxhall':                 'victoria',
    'victoria':                 'victoria',
    # W
    'warren street':            'victoria',
    'waterloo':                 'jubilee',
    'west ham':                 'district',
    'westminster':              'jubilee',
    'whitechapel':              'elizabeth',
    'wapping':                  'overground',
}


def _strip_white_bg(img_path):
    """Return a version of the PNG with white/near-white pixels made transparent (cached)."""
    cache = img_path.replace('.png', '_nobg.png')
    if os.path.exists(cache) and os.path.getsize(cache) > 200:
        return cache
    try:
        from PIL import Image as _PI
        img = _PI.open(img_path).convert('RGBA')
        pixels = img.load()
        w, h = img.size
        for y in range(h):
            for x in range(w):
                r, g, b, a = pixels[x, y]
                if r > 235 and g > 235 and b > 235:
                    pixels[x, y] = (r, g, b, 0)
        img.save(cache, 'PNG')
        return cache
    except Exception:
        return img_path


def station_badge_path(station_name):
    """Return transparent badge PNG for a station name, or None."""
    s = station_name.lower().strip()
    key = STATION_LINE_MAP.get(s)
    if not key:
        for k, v in STATION_LINE_MAP.items():
            if k in s or s in k:
                key = v
                break
    if not key:
        for keyword, k in TUBE_LINE_KEYS.items():
            if keyword in s:
                key = k
                break
    if key:
        p = os.path.join(TUBE_DIR, f'{key}.png')
        if os.path.isfile(p):
            return _strip_white_bg(p)
    return None


# Tube line keyword → badge image key
TUBE_LINE_KEYS = {
    'jubilee':          'jubilee',
    'central':          'central',
    'northern':         'northern',
    'piccadilly':       'piccadilly',
    'victoria':         'victoria',
    'district':         'district',
    'circle':           'circle',
    'bakerloo':         'bakerloo',
    'elizabeth':        'elizabeth',
    'crossrail':        'elizabeth',
    'metropolitan':     'metropolitan',
    'hammersmith & city':'hammersmith',
    'hammersmith':      'hammersmith',
    'h&c':              'hammersmith',
    'waterloo & city':  'waterloo',
    'waterloo':         'waterloo',
    'dlr':              'dlr',
    'docklands':        'dlr',
    'overground':       'overground',
    'national rail':    'overground',
    'thameslink':       'overground',
}

# ── Canvas & palette ────────────────────────────────────────────────────────────
W, H = 12192000, 6858000
M    = 550000          # left / right margin

BLUE  = '#1E22AA'
MINT  = '#ADDFB3'
BLUSH = '#F4BEAA'
AQUA  = '#96E6DC'
WHITE = '#FFFFFF'
DARK  = '#1D1D1B'
GREY  = '#6B7280'
LGREY = '#F3F4F6'
DGREY = '#9CA3AF'


# ── Primitive draw-command helpers ──────────────────────────────────────────────

def R(x, y, w, h, fill):
    return {'cmd': 'rect', 'x': x, 'y': y, 'w': w, 'h': h, 'fill': fill}


def T(text, x, y, w, h, size, color,
      font='Mont', bold=False, italic=False, align='L', wrap=True):
    return {
        'cmd':    'text',
        'text':   str(text or ''),
        'x': x, 'y': y, 'w': w, 'h': h,
        'font':   font,
        'size':   size,
        'color':  color,
        'bold':   bold,
        'italic': italic,
        'align':  align,
        'wrap':   wrap,
    }


def I(path, x, y, w, h):
    return {'cmd': 'img', 'path': path, 'x': x, 'y': y, 'w': w, 'h': h}


def L(x1, y1, x2, y2, color='#E5E7EB', width=1):
    return {'cmd': 'line', 'x1': x1, 'y1': y1, 'x2': x2, 'y2': y2,
            'color': color, 'width': width}


# ── Helper utilities ────────────────────────────────────────────────────────────

AMENITY_LABELS = {
    'wifi':                    'Wi-Fi',
    'parking':                 'Parking',
    'four_wheeler_parking':    'Car Parking',
    'two_wheeler_parking':     'Bike Parking',
    'cafeteria':               'Cafeteria',
    'gym':                     'Gym',
    'meeting_rooms':           'Meeting Rooms',
    'meeting_room':            'Meeting Room',
    'phone_booths':            'Phone Booths',
    'phone_booth':             'Phone Booth',
    'reception':               'Reception',
    'printing':                'Printing',
    'lounge':                  'Lounge',
    'security':                '24/7 Security',
    'air_conditioning':        'Air Con',
    'power_backup':            'Power Backup',
    'breakout_area':           'Breakout Area',
    'terrace':                 'Terrace',
    'bike_storage':            'Bike Storage',
    'lift':                    'Lift',
    'shower':                  'Showers',
    'kitchen':                 'Kitchen',
    'event_space':             'Event Space',
    'rooftop':                 'Rooftop',
    'onsite_cafe':             'Café',
    'onsite_cafeteria':        'Café',
    'cctv':                    'CCTV',
    'fire_safety':             'Fire Safety',
    'housekeeping':            'Housekeeping',
    'it_support':              'IT Support',
    'mail_handling':           'Mail Handling',
    'ramp_access':             'Ramp Access',
    'disabled_access':         'Disabled Access',
    'generator_backup':        'Generator',
    'storage':                 'Storage',
    'lockers':                 'Lockers',
    'nursing_room':            'Nursing Room',
    'pantry':                  'Pantry',
    'concierge':               'Concierge',
    'security_personnel':      'Security',
    'security_guard':          'Security',
    '24_7_access':             '24/7 Access',
    '24x7_access':             '24/7 Access',
    '24/7_access':             '24/7 Access',
    '24/7 access':             '24/7 Access',
    'breakout_room':           'Breakout Room',
    'recreation_room':         'Recreation',
    'lounge_area':             'Lounge Area',
    'ac':                      'Air Con',
    'power':                   'Power Backup',
    'wifi_high_speed':         'Wi-Fi',
    'high_speed_wifi':         'Wi-Fi',
    'dedicated_desks':         'Ded. Desks',
    'private_office':          'Private Office',
    'hot_desk':                'Hot Desks',
    'video_conferencing':      'Video Conf.',
    'outdoor_space':           'Outdoor Space',
    'wellness_room':           'Wellness Room',
    'fitness_center':          'Gym',
    'fitness_centre':          'Gym',
    'on_site_cafe':            'Café',
    'coffee':                  'Coffee',
    'natural_light':           'Natural Light',
    'standing_desks':          'Standing Desks',
    'accessible':              'Accessible',
    'wheelchair_access':       'Accessible',
}

# Dot colours cycling for amenity grid
_AMENITY_COLORS = [AQUA, MINT, '#7CB9E8', BLUSH, '#C8E6C9', '#FFF176', '#FFB74D']


def format_amenities(a):
    try:
        ams = json.loads(a) if isinstance(a, str) else (a or [])
    except Exception:
        ams = []
    labels = []
    for x in ams[:8]:
        slug = str(x).lower().strip()
        label = AMENITY_LABELS.get(slug)
        if not label:
            # Clean up unknown slugs
            label = slug.replace('_', ' ').replace('-', ' ').title()
            if len(label) > 20:
                label = label[:19] + '…'
        labels.append(label)
    return ', '.join(labels) if labels else 'On request'


def amenity_pill_cmds(amenities_raw, x, y, max_w, font):
    """3-column grid: coloured square dot + label text. Clean, no overflow."""
    try:
        ams = json.loads(amenities_raw) if isinstance(amenities_raw, str) else (amenities_raw or [])
    except Exception:
        ams = []
    if not ams:
        return [T('On request', x, y, max_w, 240000, 8, GREY, font=font)]

    cmds = []
    cols = 3
    col_w = max_w // cols
    row_h = 300000
    dot_size = 110000
    dot_gap  = 80000
    text_h   = 220000   # approximate height for 8pt text

    for i, slug in enumerate(ams[:9]):
        slug = str(slug).lower().strip()
        label = AMENITY_LABELS.get(slug)
        if not label:
            label = slug.replace('_', ' ').replace('-', ' ').title()
            if len(label) > 20:
                label = label[:19] + '…'

        col = i % cols
        row = i // cols
        ix = x + col * col_w
        iy = y + row * row_h

        # Vertically center both dot and text within the row
        dot_y  = iy + (row_h - dot_size) // 2
        text_y = iy + (row_h - text_h)  // 2

        color = _AMENITY_COLORS[i % len(_AMENITY_COLORS)]
        cmds.append(R(ix, dot_y, dot_size, dot_size, color))
        cmds.append(T(label,
                      ix + dot_size + dot_gap, text_y,
                      col_w - dot_size - dot_gap - 40000, text_h,
                      8, DARK, font=font, wrap=False))

    return cmds


def price_str(c):
    p = c.get('price_from') or c.get('price_str')
    if not p:
        return 'On request'
    if isinstance(p, str):
        return p
    unit = (c.get('price_unit') or 'MONTHLY').upper()
    suffix = {'MONTHLY': '/mo', 'DAILY': '/day', 'HOURLY': '/hr'}.get(unit, '/mo')
    city = (c.get('city') or '').lower()
    symbol = '£' if 'london' in city or 'uk' in city else '₹'
    return f'{symbol}{int(p):,}{suffix}/seat'


def extract_transport_parts(transport_str):
    """Return (tube_line, bus_line, other) from a transport string."""
    s = transport_str or ''
    tube = bus = other = ''
    for part in s.split(','):
        part = part.strip()
        pl = part.lower()
        if 'tube' in pl or 'underground' in pl or 'metro' in pl or 'line' in pl:
            tube = part
        elif 'bus' in pl:
            bus = part
        elif not other:
            other = part
    return tube, bus, other


def tube_badge_path(transport_str, raw=False):
    """Return tube line badge PNG path for transport_str, or None.
    raw=True returns the original PNG (white background intact — use on white/light backgrounds).
    raw=False (default) strips white bg for use on coloured/dark backgrounds.
    """
    s = (transport_str or '').lower()
    station = _extract_station_name(transport_str)
    if station:
        s2 = station.lower().strip()
        key = STATION_LINE_MAP.get(s2)
        if not key:
            for k, v in STATION_LINE_MAP.items():
                if k in s2 or s2 in k:
                    key = v
                    break
        if key:
            p = os.path.join(TUBE_DIR, f'{key}.png')
            if os.path.isfile(p):
                return p if raw else _strip_white_bg(p)
    for keyword, key in TUBE_LINE_KEYS.items():
        if keyword in s:
            p = os.path.join(TUBE_DIR, f'{key}.png')
            if os.path.isfile(p):
                return p if raw else _strip_white_bg(p)
    return None


def get_logo_png(white=False):
    """Return a path to a usable PNG of the myHQ logo (transparent background).
    white=True returns the all-white version for use on dark/coloured backgrounds.
    """
    suffix = '-white' if white else ''
    png_path = os.path.join(BASE_DIR, 'static', 'images', f'myhq-logo{suffix}.png')
    if os.path.exists(png_path) and os.path.getsize(png_path) > 5000:
        return png_path

    # If white version requested but missing, generate from the colour version
    if white:
        base = get_logo_png(white=False)
        try:
            import numpy as np
            from PIL import Image as _PI
            img = _PI.open(base).convert('RGBA')
            data = np.array(img)
            visible = data[:, :, 3] > 10
            data[visible, 0] = 255
            data[visible, 1] = 255
            data[visible, 2] = 255
            _PI.fromarray(data, 'RGBA').save(png_path, 'PNG')
            return png_path
        except Exception:
            return base   # fall back to colour logo

    os.makedirs(os.path.dirname(png_path), exist_ok=True)

    # Candidate SVG sources
    svg_candidates = [
        os.path.join(BASE_DIR, 'static', 'images', 'myhq-logo.svg'),
        os.path.join(BASE_DIR, '..', 'onboarding-tool', 'static', 'myhq-logo.svg'),
        os.path.join(BASE_DIR, 'static', 'myhq-logo.svg'),
    ]
    logo_svg = next((p for p in svg_candidates if os.path.exists(p)), None)

    if logo_svg:
        # Try qlmanage (macOS built-in) and strip white background
        try:
            import numpy as np
            from PIL import Image as _PI
            tmp_dir = os.path.join(BASE_DIR, 'static', 'images')
            subprocess.run(['qlmanage', '-t', '-s', '1280', '-o', tmp_dir, logo_svg],
                           timeout=15, capture_output=True)
            ql_candidate = os.path.join(tmp_dir, os.path.basename(logo_svg) + '.png')
            if os.path.exists(ql_candidate):
                img = _PI.open(ql_candidate).convert('RGBA')
                data = np.array(img)
                white_mask = (data[:,:,0] > 240) & (data[:,:,1] > 240) & (data[:,:,2] > 240)
                data[white_mask, 3] = 0
                result = _PI.fromarray(data, 'RGBA')
                bbox = result.getbbox()
                if bbox:
                    result = result.crop(bbox)
                result.save(png_path, 'PNG')
                os.remove(ql_candidate)
                if os.path.getsize(png_path) > 5000:
                    return png_path
        except Exception:
            pass

    # Fallback: draw "myHQ" text with Pillow
    try:
        from PIL import Image as PILImage, ImageDraw, ImageFont
        img = PILImage.new('RGBA', (320, 120), (255, 255, 255, 0))
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 64)
        except Exception:
            font = ImageFont.load_default()
        draw.text((10, 20), 'myHQ', fill=(30, 34, 170, 255), font=font)
        img.save(png_path, 'PNG')
        return png_path
    except Exception:
        pass

    try:
        from PIL import Image as PILImage
        PILImage.new('RGBA', (1, 1), (255, 255, 255, 0)).save(png_path, 'PNG')
    except Exception:
        pass
    return png_path


# ── Logo helpers ────────────────────────────────────────────────────────────────

def _logo_bg_brightness(image_path):
    """Return average brightness (0-255) of the image region behind the top-right logo."""
    try:
        from PIL import Image as _PI
        img = _PI.open(image_path).convert('RGB')
        iw, ih = img.size
        # Logo sits at top-right: x = W-1200000, y=160000, w=1100000, h=412000 (EMU)
        px = max(0, int((W - 1200000) / W * iw))
        py = max(0, int(160000 / H * ih))
        pw = max(1, int(1100000 / W * iw))
        ph = max(1, int(412000 / H * ih))
        px = min(px, iw - pw)
        py = min(py, ih - ph)
        region = img.crop((px, py, px + pw, py + ph))
        pixels = list(region.getdata())
        return sum(0.299 * r + 0.587 * g + 0.114 * b for r, g, b in pixels) / max(len(pixels), 1)
    except Exception:
        return 255  # assume light → colored logo


def _pick_logo_path(bg_hex=None, bg_image=None):
    """Return logo path — white variant on dark bg, colored on light bg.
    Both logos already have transparent backgrounds so no stripping needed."""
    use_white = False
    if bg_image and os.path.isfile(str(bg_image)):
        try:
            use_white = _logo_bg_brightness(bg_image) < 140
        except Exception:
            pass
    elif bg_hex:
        h = bg_hex.lstrip('#')
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        use_white = (0.299 * r + 0.587 * g + 0.114 * b) < 128
    return get_logo_png(white=use_white)


# ── Shared slide builders (used by both templates) ─────────────────────────────

_LOGO_X, _LOGO_Y, _LOGO_W, _LOGO_H = W - 1200000, 140000, 1100000, 412000


def _small_logo(logo_path, bg='#FFFFFF', bg_image=None):
    """Top-right logo — transparent, white or colored chosen by background brightness."""
    path = _pick_logo_path(bg_hex=bg, bg_image=bg_image)
    return I(path, _LOGO_X, _LOGO_Y, _LOGO_W, _LOGO_H)


def _small_logo_transparent(logo_path):
    """Top-right logo over map/coloured background — white-stripped transparent PNG."""
    return [I(_strip_white_bg(logo_path), _LOGO_X, _LOGO_Y, _LOGO_W, _LOGO_H)]


def _comparison_table_slide(all_centres, header_fill, header_text_color, font, logo):
    """Slide 2: Space Comparison table with photo column. Vertical layout when > 5 centres."""
    sl = [R(0, 0, W, H, WHITE)]
    sl.append(_small_logo(logo))

    sl.append(T('Space Comparison', M, 250000, W - 2 * M, 450000,
                22, BLUE, font=font, bold=True))
    sl.append(T('All shortlisted workspaces at a glance', M, 720000, W - 2 * M, 300000,
                11, GREY, font=font))

    table_x = M
    table_y = 1100000
    table_w = W - 2 * M
    table_h = H - 1500000

    num_centres = max(len(all_centres), 1)
    hdr_h   = 520000
    raw_row_h = (table_h - hdr_h) // num_centres
    row_h   = max(400000, min(1100000, raw_row_h))

    # Use a photo column – 900000 wide if rows are tall enough
    photo_col_w = 900000 if row_h >= 600000 else 0

    # Column definitions: (label, fixed_width_or_None, field_fn)
    def _centre_tube_text(c):
        """Station name only — no badge image in comparison table."""
        transport = c.get('transport', '') or ''
        station = _extract_station_name(transport)
        return station or '—'

    def _centre_name(c):
        name = c.get('name', '—')
        brand = (c.get('brand') or '').strip()
        # Prepend brand if it's not already in the name
        if brand and brand.lower() not in name.lower():
            return f'{brand} – {name}'
        return name

    def _centre_hours(c):
        h = c.get('open_hours') or ''
        if h and h.strip():
            return h.strip()
        return '9:00 AM – 6:00 PM'

    def _centre_type(c):
        t = c.get('space_type') or ''
        if not t:
            # Try to infer from pricing type for London map entries
            try:
                import json as _j
                pricing = c.get('pricing') or []
                if isinstance(pricing, str):
                    pricing = _j.loads(pricing)
                if pricing and isinstance(pricing, list):
                    t = pricing[0].get('type', '') or ''
            except Exception:
                pass
        return t.replace('_', ' ').title() if t else '—'

    fixed_cols = [
        ('#',           380000,  lambda i, c: str(i + 1)),
        ('NAME',       2100000,  lambda i, c: _centre_name(c)),
        ('ADDRESS',    2000000,  lambda i, c: c.get('address', '—')),
        ('PRICE/SEAT', 1200000,  lambda i, c: price_str(c)),
        ('NEAREST TUBE', 1700000, lambda i, c: _centre_tube_text(c)),
        ('HOURS',      1300000,  lambda i, c: _centre_hours(c)),
        ('TYPE',       1012000,  lambda i, c: _centre_type(c)),
    ]
    if photo_col_w:
        fixed_cols.append(('PHOTO', photo_col_w, None))

    # Draw header
    cx = table_x
    for label, cw, _ in fixed_cols:
        sl.append(R(cx, table_y, cw, hdr_h, header_fill))
        sl.append(T(label, cx + 35000, table_y + 110000, cw - 70000, hdr_h - 180000,
                    8, header_text_color, font=font, bold=True))
        cx += cw

    # Draw data rows
    for ri, centre in enumerate(all_centres):
        ry = table_y + hdr_h + ri * row_h
        row_fill = WHITE if ri % 2 == 0 else LGREY
        cx = table_x
        for label, cw, val_fn in fixed_cols:
            sl.append(R(cx, ry, cw, row_h, row_fill))
            if val_fn is None:
                # Photo cell
                imgs = centre.get('images', [])
                photo = imgs[0] if imgs else None
                sl.append(I(photo, cx + 20000, ry + 20000,
                            cw - 40000, row_h - 40000))
            else:
                val = val_fn(ri, centre)
                sl.append(T(val, cx + 35000, ry + 60000,
                            cw - 70000, row_h - 100000,
                            8, DARK, font=font, wrap=True))
            cx += cw
        sl.append(L(table_x, ry + row_h, table_x + table_w, ry + row_h,
                    color='#E5E7EB', width=1))

    return sl


def _client_requirements_slide(proposal, template, db_centres, manual_centres,
                                font, logo,
                                label_fill, label_text, date_label_color):
    """Slide 4: Client Requirements."""
    sl = [R(0, 0, W, H, WHITE)]
    sl.append(_small_logo(logo))
    sl.append(T('Client Requirements', M, 250000, 6000000, 400000,
                22, BLUE, font=font, bold=True))

    today = datetime.date.today().strftime('%d %B %Y')

    reqs = [
        ('CLIENT',         proposal.get('client_name')     or '—'),
        ('COMPANY',        proposal.get('client_company')  or '—'),
        ('LOCATION',       proposal.get('client_location') or '—'),
        ('TEAM SIZE',      proposal.get('team_size')       or '—'),
        ('SPACE TYPE',     proposal.get('space_type')      or '—'),
        ('BUDGET',         proposal.get('budget')          or '—'),
    ]

    lbl_col_w  = 2100000
    val_col_w  = 4700000
    row_h      = 660000
    tbl_x      = M
    tbl_y      = 1000000

    for i, (lbl, val) in enumerate(reqs):
        ry   = tbl_y + i * row_h
        vfill = WHITE if i % 2 == 0 else LGREY
        sl.append(R(tbl_x,               ry, lbl_col_w, row_h, label_fill))
        sl.append(R(tbl_x + lbl_col_w,   ry, val_col_w, row_h, vfill))
        sl.append(T(lbl, tbl_x + 60000, ry + 180000, lbl_col_w - 120000, row_h - 200000,
                    9, label_text, font=font, bold=True))
        sl.append(T(val, tbl_x + lbl_col_w + 80000, ry + 150000,
                    val_col_w - 160000, row_h - 200000,
                    11, DARK, font=font, wrap=True))
        # bottom divider
        sl.append(L(tbl_x, ry + row_h,
                    tbl_x + lbl_col_w + val_col_w, ry + row_h))

    # Right summary card – shows proposal meta only, no repeated fields
    card_x = tbl_x + lbl_col_w + val_col_w + 300000
    card_w = W - card_x - M
    card_h = len(reqs) * row_h
    card_y = tbl_y

    sl.append(R(card_x, card_y, card_w, card_h, BLUE))

    inner_x = card_x + 120000
    inner_w  = card_w - 240000

    # Date label
    sl.append(T(today, inner_x, card_y + 120000, inner_w, 250000,
                9, date_label_color, font=font))
    sl.append(L(card_x + 60000, card_y + 400000,
                card_x + card_w - 60000, card_y + 400000,
                color='#3A3E99', width=1))

    # Company name (prominent heading)
    company = proposal.get('client_company') or proposal.get('client_name') or '—'
    sl.append(T(company, inner_x, card_y + 460000, inner_w, 900000,
                28, WHITE, font=font, bold=True, wrap=True))

    sl.append(L(card_x + 60000, card_y + 1480000,
                card_x + card_w - 60000, card_y + 1480000,
                color='#3A3E99', width=1))

    # "Prepared by myHQ" section
    sl.append(T('PREPARED BY', inner_x, card_y + 1560000, inner_w, 240000,
                8, date_label_color, font=font, bold=True))
    sl.append(T('myHQ by Anarock', inner_x, card_y + 1840000, inner_w, 360000,
                16, WHITE, font=font, bold=True))

    sl.append(L(card_x + 60000, card_y + 2280000,
                card_x + card_w - 60000, card_y + 2280000,
                color='#3A3E99', width=1))

    # Contact
    sl.append(T('YOUR CONSULTANT', inner_x, card_y + 2380000, inner_w, 240000,
                8, date_label_color, font=font, bold=True))
    sl.append(T('arjun.budhiraja@myhq.in', inner_x, card_y + 2660000, inner_w, 280000,
                11, WHITE, font=font))
    sl.append(T('+44 7863 754009', inner_x, card_y + 2980000, inner_w, 280000,
                11, WHITE, font=font))

    # Decorative myHQ logo area — fixed size preserving 1280×479 aspect ratio
    sl.append(R(card_x, card_y + card_h - 600000, card_w, 600000, '#16199A'))
    _logo_w = get_logo_png(white=True)
    _lw, _lh = 900000, int(900000 * 479 // 1280)   # 900000 × 337000, ratio exact → no cover-crop
    _lx = card_x + (card_w - _lw) // 2
    _ly = card_y + card_h - 600000 + (600000 - _lh) // 2
    sl.append(I(_logo_w, _lx, _ly, _lw, _lh))

    return sl


def _extract_station_name(transport_str):
    """Pull the nearest station name from a transport string."""
    s = (transport_str or '').strip()
    if not s:
        return ''
    import re

    def _clean(name):
        """Strip trailing '(Line name)' parentheticals and whitespace."""
        return re.sub(r'\s*\([^)]*\)', '', name).strip()

    # Pattern: "X mins walk from Y station", "nearest: Y", etc.
    # Use \b so "to" inside "Victoria" is not matched as a keyword
    m = re.search(r'\b(?:from|to|nearest[:\s]+)\s*([A-Z][^,\n;(]{2,30}?)(?:\s*station|\s*tube|\s*underground|\s*\(|,|$)',
                  s, re.IGNORECASE)
    if m:
        return _clean(m.group(1))

    # Fallback: first proper-noun segment, strip parentheticals
    parts = re.split(r'[,;:\n]', s)
    for p in parts:
        p = p.strip()
        p = re.sub(r'^(bus|tube|underground|overground|rail|dlr|metro)\s*:?\s*', '', p, flags=re.IGNORECASE)
        p = _clean(p)
        if len(p) > 3:
            return p[:40]
    return _clean(s[:40])


def _extract_walk_time(transport_str):
    """Pull walking time from transport string, e.g. '5 min walk'."""
    import re
    m = re.search(r'(\d+)\s*(?:min|minute|mins)', transport_str or '', re.IGNORECASE)
    if m:
        return f"{m.group(1)} min walk"
    return ''


def _centre_slide(idx, centre, template, font, logo):
    """One centre detail slide (Slide 6+)."""
    sl = [R(0, 0, W, H, WHITE)]

    # ── Accent bar at top ────────────────────────────────────────────────────────
    sl.append(R(0, 0, W, 120000, AQUA if template != 'india' else BLUE))

    # ── Header ──────────────────────────────────────────────────────────────────
    sl.append(T(f'{idx:02d}', M, 140000, 300000, 480000,
                22, BLUE, font=font, bold=True))
    sl.append(T(centre.get('name', 'Centre'),
                M + 380000, 160000, 7600000, 420000,
                18, DARK, font=font, bold=True))
    sl.append(L(M, 680000, W - M, 680000, color=LGREY, width=1))

    # ── Left column ──────────────────────────────────────────────────────────────
    lcol_x = M
    lcol_w  = 4000000
    cy = 760000

    # ADDRESS
    sl.append(T('ADDRESS', lcol_x, cy, lcol_w, 190000, 7, DGREY, font=font, bold=True))
    cy += 190000
    sl.append(T(centre.get('address') or '—', lcol_x, cy, lcol_w, 280000, 9, DARK, font=font, wrap=True))
    cy += 310000
    sl.append(L(lcol_x, cy, lcol_x + lcol_w, cy, color='#EEEEEE', width=1))
    cy += 50000

    # PRICE + HOURS on same row
    half = (lcol_w - 80000) // 2
    sl.append(T('PRICE / SEAT', lcol_x, cy, half, 190000, 7, DGREY, font=font, bold=True))
    sl.append(T('OPEN HOURS', lcol_x + half + 80000, cy, half, 190000, 7, DGREY, font=font, bold=True))
    cy += 190000
    sl.append(T(price_str(centre), lcol_x, cy, half, 280000, 10, BLUE, font=font, bold=True))
    sl.append(T(centre.get('open_hours') or '9:00 AM – 6:00 PM',
                lcol_x + half + 80000, cy, half, 280000, 9, DARK, font=font))
    cy += 340000
    sl.append(L(lcol_x, cy, lcol_x + lcol_w, cy, color='#EEEEEE', width=1))
    cy += 60000

    # AMENITIES as pills
    sl.append(T('AMENITIES', lcol_x, cy, lcol_w, 190000, 7, DGREY, font=font, bold=True))
    cy += 200000
    amenity_cmds = amenity_pill_cmds(centre.get('amenities', '[]'), lcol_x, cy, lcol_w, font)
    sl.extend(amenity_cmds)
    # Grid is 3-column; calculate actual rows used
    try:
        ams = json.loads(centre.get('amenities', '[]') if isinstance(centre.get('amenities', '[]'), str) else '[]')
        n_ams = min(len(ams), 9)
    except Exception:
        n_ams = 0
    grid_rows = max(1, (n_ams + 2) // 3)
    cy += grid_rows * 290000 + 100000

    sl.append(L(lcol_x, cy, lcol_x + lcol_w, cy, color='#EEEEEE', width=1))
    cy += 60000

    # NEAREST TUBE / TRANSPORT
    transport_str = centre.get('transport', '')
    tube, bus, other = extract_transport_parts(transport_str)
    badge = tube_badge_path(transport_str)
    station_name = _extract_station_name(transport_str)
    walk_time = _extract_walk_time(transport_str)

    sl.append(T('NEAREST TUBE / TRANSPORT', lcol_x, cy, lcol_w, 190000, 7, DGREY, font=font, bold=True))
    cy += 200000

    if badge:
        badge_h = 210000   # 700000 / 3.33 = matches 320×96 badge ratio exactly
        sl.append(I(badge, lcol_x, cy, 700000, badge_h))
        txt_x = lcol_x + 760000
        txt_w = lcol_w - 760000
        if station_name:
            sl.append(T(station_name, txt_x, cy, txt_w, 200000, 9, DARK, font=font, bold=True))
        if walk_time:
            sl.append(T(walk_time, txt_x, cy + 190000, txt_w, 180000, 8, GREY, font=font))
        cy += 300000
    elif tube or station_name or other:
        display = station_name or tube or other
        sl.append(R(lcol_x, cy + 70000, 8000, 140000, BLUE))
        sl.append(T(display, lcol_x + 60000, cy, lcol_w - 60000, 200000, 9, DARK, font=font, bold=True))
        if walk_time:
            sl.append(T(walk_time, lcol_x + 60000, cy + 210000, lcol_w - 60000, 180000, 8, GREY, font=font))
        cy += 320000

    if bus:
        sl.append(R(lcol_x, cy + 70000, 8000, 140000, AQUA if template != 'india' else BLUE))
        sl.append(T(bus, lcol_x + 60000, cy, lcol_w - 60000, 200000, 9, DARK, font=font))
        cy += 270000

    cy += 60000
    sl.append(L(lcol_x, cy, lcol_x + lcol_w, cy, color='#EEEEEE', width=1))
    cy += 60000

    # WHY WE RECOMMEND
    why_text = (centre.get('why_recommend') or centre.get('about')
                or 'A premium flexible workspace in a prime location, ideal for growing teams.')
    sl.append(T('WHY WE RECOMMEND', lcol_x, cy, lcol_w, 190000, 7, BLUE, font=font, bold=True))
    cy += 200000
    # Fill remaining left-column height
    remaining = H - cy - 100000
    sl.append(T(why_text, lcol_x, cy, lcol_w, max(remaining, 400000), 9, DARK, font=font, wrap=True))

    # ── Right column: 2×2 equal photo grid ───────────────────────────────────────
    GAP  = 120000
    rx   = M + lcol_w + 200000
    ry   = 720000
    rw   = W - rx - 80000
    rh   = H - ry - 100000

    cw   = (rw - GAP) // 2
    ch   = (rh - GAP) // 2

    images = centre.get('images', [])

    def _cell(ci, gx, gy):
        p = images[ci] if ci < len(images) else None
        return I(p, gx, gy, cw, ch)

    sl.append(_cell(0, rx,            ry))
    sl.append(_cell(1, rx + cw + GAP, ry))
    sl.append(_cell(2, rx,            ry + ch + GAP))
    sl.append(_cell(3, rx + cw + GAP, ry + ch + GAP))

    # Logo last so it sits on top — pick white/colored based on top-right image brightness
    top_right = images[1] if len(images) > 1 else (images[0] if images else None)
    sl.append(_small_logo(logo, bg_image=top_right))

    return sl


def _existing_clients_slide(font, logo):
    # Full-slide image taken directly from the reference PDF
    slide_img = os.path.join(IMG_DIR, 'slide_existing_clients.png')
    return [I(slide_img, 0, 0, W, H)]


def _testimonials_slide(font, logo, bg_color=WHITE, heading_color=DARK):
    # Full-slide image taken directly from the reference PDF
    slide_img = os.path.join(IMG_DIR, 'slide_word_of_mouth.png')
    return [I(slide_img, 0, 0, W, H)]


def generate_proposal_map(centres, out_path):
    """Generate a static map image showing the selected centres with numbered pins.
    Uses OSM tiles + PIL only — no browser/Playwright required."""
    import math, urllib.request, io as _io
    from PIL import Image as _PI, ImageDraw as _ID, ImageFont as _IF

    # ── 1. Parse coordinates ────────────────────────────────────────────────
    points = []
    for c in centres:
        lat = lng = None
        # Direct lat/lng fields (manual entries from London map)
        if c.get('lat') and c.get('lng'):
            try:
                lat, lng = float(c['lat']), float(c['lng'])
            except Exception:
                pass
        # Fallback: 'coordinates' field in 'lng;lat' format (DB centres)
        if lat is None:
            coords = c.get('coordinates', '')
            if coords and ';' in coords:
                try:
                    lng_s, lat_s = coords.split(';')
                    lat, lng = float(lat_s), float(lng_s)
                except Exception:
                    pass
        if lat is not None and lng is not None:
            # Try swapping if the values look backwards (lat should be ~51, lng ~-0.x)
            if not (51.0 <= lat <= 52.0) and (51.0 <= lng <= 52.0):
                lat, lng = lng, lat
            # Filter to Greater London bounds
            if 51.2 <= lat <= 51.8 and -0.6 <= lng <= 0.4:
                points.append({'lat': lat, 'lng': lng, 'name': c.get('name', '')})
    if not points:
        return None

    # ── 2. Tile helpers ─────────────────────────────────────────────────────
    def _deg2tile(lat, lng, z):
        n = 2 ** z
        xt = (lng + 180) / 360 * n
        yt = (1 - math.log(math.tan(math.radians(lat)) + 1 / math.cos(math.radians(lat))) / math.pi) / 2 * n
        return xt, yt

    def _pick_zoom(points, img_w, img_h, pad=100):
        """Pick highest zoom where all points fit inside the image with padding."""
        for z in range(15, 8, -1):  # try down to z=9 (covers all of London)
            xs = [_deg2tile(p['lat'], p['lng'], z)[0] for p in points]
            ys = [_deg2tile(p['lat'], p['lng'], z)[1] for p in points]
            span_x = (max(xs) - min(xs)) * TILE
            span_y = (max(ys) - min(ys)) * TILE
            if span_x <= img_w - pad * 2 and span_y <= img_h - pad * 2:
                return z
        return 9  # zoom 9 comfortably shows all of Greater London

    # Render at 2x then scale down — gives crisp anti-aliased output.
    # Map panel on the slide is nearly square (ratio ≈ 1.0).
    # Generate at matching ratio so cover-crop removes almost nothing.
    IMG_W, IMG_H = 960, 960
    SCALE = 2
    RW, RH = IMG_W * SCALE, IMG_H * SCALE
    TILE = 512  # CartoDB @2x tiles are 512×512

    zoom = _pick_zoom(points, RW, RH) if len(points) > 1 else 14

    clat = sum(p['lat'] for p in points) / len(points)
    clng = sum(p['lng'] for p in points) / len(points)
    cx, cy = _deg2tile(clat, clng, zoom)

    tiles_x = math.ceil(RW / TILE) + 4   # extra buffer so edge points aren't cut
    tiles_y = math.ceil(RH / TILE) + 4
    tx0 = int(cx) - tiles_x // 2
    ty0 = int(cy) - tiles_y // 2

    # ── 3. Fetch & stitch CartoDB Voyager @2x tiles ─────────────────────────
    import ssl as _ssl
    _ctx = _ssl._create_unverified_context()
    canvas = _PI.new('RGB', ((tiles_x + 1) * TILE, (tiles_y + 1) * TILE), (242, 243, 244))
    headers = {'User-Agent': 'myHQ-proposal-tool/1.0'}
    n_tiles = 2 ** zoom
    subdomain = ['a', 'b', 'c', 'd']
    for dx in range(tiles_x + 1):
        for dy in range(tiles_y + 1):
            tx = (tx0 + dx) % n_tiles
            ty = (ty0 + dy) % n_tiles
            if ty < 0 or ty >= n_tiles:
                continue
            s = subdomain[(tx + ty) % 4]
            url = f'https://{s}.basemaps.cartocdn.com/rastertiles/voyager/{zoom}/{tx}/{ty}@2x.png'
            try:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, timeout=8, context=_ctx) as r:
                    tile = _PI.open(_io.BytesIO(r.read())).convert('RGB')
                    if tile.size != (TILE, TILE):
                        tile = tile.resize((TILE, TILE), _PI.LANCZOS)
                canvas.paste(tile, (dx * TILE, dy * TILE))
            except Exception:
                pass

    # ── 4. Crop to render size centred on points ────────────────────────────
    px_cx = (cx - tx0) * TILE
    px_cy = (cy - ty0) * TILE
    left  = int(px_cx - RW / 2)
    top   = int(px_cy - RH / 2)
    img   = canvas.crop((left, top, left + RW, top + RH))

    # ── 5. Draw markers — all pixel values at 2x (canvas is 2x, downscaled later) ──
    draw = _ID.Draw(img)
    # Font sizes at 2x so they appear at correct size after downscale to IMG_W x IMG_H
    fs_label, fs_num = 30, 40

    def _load_font(size):
        """Load a bold font at the given size, trying multiple paths."""
        candidates = [
            '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',        # Debian/Ubuntu
            '/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf',
            '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf', # Ubuntu
            '/usr/share/fonts/liberation/LiberationSans-Bold.ttf',
            '/usr/share/fonts/truetype/freefont/FreeSansBold.ttf',
            '/System/Library/Fonts/Helvetica.ttc',                          # macOS
            '/Library/Fonts/Arial Bold.ttf',
        ]
        for path in candidates:
            try:
                return _IF.truetype(path, size)
            except Exception:
                pass
        # Pillow 10.1+ supports size on load_default (uses a bundled TrueType font)
        try:
            return _IF.load_default(size=size)
        except Exception:
            return _IF.load_default()

    font_label = _load_font(fs_label)
    font_num   = _load_font(fs_num)

    for i, p in enumerate(points, 1):
        px, py = _deg2tile(p['lat'], p['lng'], zoom)
        px = int((px - tx0) * TILE) - left
        py = int((py - ty0) * TILE) - top

        # Drop-shadow (2x offsets)
        r = 56
        draw.ellipse([px - r + 6, py - r + 6, px + r + 6, py + r + 6],
                     fill=(0, 0, 0, 60) if img.mode == 'RGBA' else '#cccccc')
        # Blue circle marker
        draw.ellipse([px - r, py - r, px + r, py + r], fill='#1E22AA', outline='white', width=6)
        num_text = str(i)
        bb = draw.textbbox((0, 0), num_text, font=font_num)
        tw, th = bb[2] - bb[0], bb[3] - bb[1]
        draw.text((px - tw // 2, py - th // 2 - 2), num_text, fill='white', font=font_num)

        # Label pill — show full name, shift left if it would clip the right edge
        label = p['name']
        lb = draw.textbbox((0, 0), label, font=font_label)
        lw, lh = lb[2] - lb[0], lb[3] - lb[1]
        pad = 20
        lx = px + r + 20
        ly = py - lh // 2 - pad
        # Prefer right of pin; flip left if it clips; truncate with ellipsis if still too wide
        if lx + lw + pad * 2 > RW - 20:
            lx = px - r - lw - pad * 2 - 20
        if lx < 10:
            max_w = RW - (px + r + 20) - pad * 2 - 20
            while lw > max_w and len(label) > 6:
                label = label[:-2] + '…'
                lb = draw.textbbox((0, 0), label, font=font_label)
                lw, lh = lb[2] - lb[0], lb[3] - lb[1]
            lx = px + r + 20
        draw.rounded_rectangle([lx, ly, lx + lw + pad * 2, ly + lh + pad * 2],
                                radius=12, fill='white', outline='#d1d5db', width=4)
        draw.text((lx + pad, ly + pad), label, fill='#1E22AA', font=font_label)

    # ── 6. Scale down to final size (anti-aliased) ──────────────────────────
    img = img.resize((IMG_W, IMG_H), _PI.LANCZOS)
    img.save(out_path, 'PNG', optimize=True)
    return out_path


_LOCATION_DATA = {
    # key: lowercase keywords to match → (about_text, highlights, [popular_stations])
    'city of london': (
        'The City of London is the UK\'s premier financial district — a square mile packed with global banks, law firms, and tech companies. With Crossrail, Tube, and rail all converging here, getting in and out is effortless.',
        ['One of the world\'s leading financial centres', 'Crossrail (Elizabeth line) cuts journey times dramatically', 'Multiple Tube lines: Central, Circle, District, Jubilee', 'World-class dining and networking at every corner'],
        ['Liverpool Street', 'Bank', 'Cannon Street', 'Moorgate', 'Aldgate'],
    ),
    'canary wharf': (
        'Canary Wharf is London\'s second financial hub, home to HSBC, Barclays, and hundreds of leading firms. Superbly connected via the Jubilee line, DLR, and Elizabeth line, it offers modern workspaces in a landmark riverside setting.',
        ['Home to major global banks and financial institutions', 'Elizabeth line direct to Heathrow & Paddington', 'Jubilee line to Westminster & London Bridge in under 10 mins', 'Extensive shopping, dining, and leisure on-site'],
        ['Canary Wharf', 'Heron Quays', 'South Quay', 'Crossharbour'],
    ),
    'shoreditch': (
        'Shoreditch is London\'s creative and tech heartland, attracting startups, agencies, and fast-growing scaleups. The area buzzes with innovation, with a tight-knit business community and superb east London connections.',
        ['London\'s leading tech and creative hub', 'Excellent connections via Overground and Elizabeth line', 'Walking distance to Liverpool Street and Old Street stations', 'Vibrant independent food, coffee, and social scene'],
        ['Shoreditch High Street', 'Old Street', 'Liverpool Street', 'Bethnal Green'],
    ),
    'old street': (
        'Old Street — the heart of Silicon Roundabout — is where London\'s tech ecosystem thrives. Surrounded by co-working pioneers, accelerators, and VC-backed startups, this is the ideal base for forward-thinking businesses.',
        ['Centre of London\'s tech startup ecosystem', 'Northern line direct to King\'s Cross and the West End', 'Huge choice of flexible and serviced offices nearby', 'Bike-friendly streets and excellent bus connections'],
        ['Old Street', 'Moorgate', 'Angel', 'Barbican'],
    ),
    'mayfair': (
        'Mayfair is synonymous with prestige — London\'s most exclusive business address. Home to luxury brands, private equity, hedge funds, and embassies, offices here send an unmistakable signal of quality and ambition.',
        ['London\'s most prestigious business postcode', 'Steps from Bond Street and Green Park stations', 'Surrounded by five-star hotels and fine dining', 'Green Park and Hyde Park nearby for client walks'],
        ['Bond Street', 'Green Park', 'Oxford Circus', 'Marble Arch'],
    ),
    'victoria': (
        'Victoria is one of London\'s best-connected transport hubs, giving easy access to Heathrow, Gatwick, and the whole of the Underground network. It is a cost-effective alternative to Mayfair with a rapidly improving offer of quality workspaces.',
        ['Direct trains to Gatwick Airport in under 30 minutes', 'Victoria, Circle, and District lines converge here', 'Strong presence of government, NGOs, and professional services', 'Major regeneration driving new workspace supply'],
        ['Victoria', 'Sloane Square', 'St James\'s Park', 'Pimlico'],
    ),
    'paddington': (
        'Paddington has been transformed by the Elizabeth line into one of London\'s fastest-growing business districts. With direct links to Heathrow in 15 minutes and the City in under 10, it offers unmatched connectivity for client-facing teams.',
        ['Elizabeth line to Heathrow Airport in 15 minutes', 'Direct Heathrow Express rail link on-site', 'Rapidly growing cluster of tech and media firms', 'Canal-side setting with excellent work-life amenities'],
        ['Paddington', 'Edgware Road', 'Bayswater', 'Royal Oak'],
    ),
    'hammersmith': (
        'Hammersmith is a well-established west London business district, popular with media, broadcasting, and professional services firms. Excellent Tube and Overground connections make it ideal for teams working across the capital.',
        ['Major media and broadcasting hub (BBC, Disney, L\'Oreal)', 'Four Tube lines: Piccadilly, District, Circle, Hammersmith & City', 'Riverside setting with excellent restaurants and bars', 'Competitive pricing versus central London'],
        ['Hammersmith', 'Ravenscourt Park', 'Stamford Brook', 'Barons Court'],
    ),
    'waterloo': (
        'Waterloo sits at the crossroads of south and central London. The largest rail terminus in the UK ensures unrivalled connectivity, while the South Bank\'s cultural energy makes it a dynamic and inspiring place to work.',
        ['Busiest rail terminus in the UK', 'Jubilee and Bakerloo lines plus overground services', 'South Bank cultural quarter: Tate Modern, Southbank Centre', 'Fast river boat services to Canary Wharf and the City'],
        ['Waterloo', 'Lambeth North', 'Southwark', 'London Bridge'],
    ),
    'london bridge': (
        'London Bridge bridges (literally) the City and South Bank, offering a vibrant mix of workspace types right next to one of London\'s busiest transport hubs. The Borough Market area adds to its appeal with world-class food and culture.',
        ['Seconds from London Bridge Station (Jubilee & Northern lines)', 'Walking distance to the City\'s financial core', 'Borough Market and Bermondsey Street for client entertainment', 'Strong legal, financial, and tech cluster'],
        ['London Bridge', 'Borough', 'Bermondsey', 'Elephant & Castle'],
    ),
    'farringdon': (
        'Farringdon is rapidly becoming one of London\'s most sought-after business locations, thanks to Crossrail making it a new interchange supernode. The area blends the creative energy of Clerkenwell with excellent connectivity.',
        ['New Elizabeth line interchange — connections in every direction', 'Heart of London\'s legal and media village', 'Clerkenwell\'s thriving design and architecture scene nearby', 'Excellent cycling infrastructure'],
        ['Farringdon', 'Barbican', 'Chancery Lane', 'Clerkenwell'],
    ),
    'king\'s cross': (
        'King\'s Cross has been completely reimagined over the past decade and is now home to Google, Facebook, and a thriving knowledge economy cluster. Six Tube lines, Eurostar, and intercity rail make it the most connected location in London.',
        ['Six Tube lines converge at King\'s Cross St. Pancras', 'Eurostar terminal for Paris/Brussels in under 2.5 hours', 'Google UK headquarters and growing tech campus', 'Coal Drops Yard — exceptional dining and retail'],
        ['King\'s Cross St Pancras', 'Euston', 'Russell Square', 'Caledonian Road'],
    ),
    'euston': (
        'Euston is a major transport gateway for those travelling to and from the Midlands and the North. With HS2 investment transforming the area, Euston is set to become one of London\'s most important business locations.',
        ['Major terminus for trains to Birmingham, Manchester, and Scotland', 'Victoria and Northern lines on the doorstep', 'HS2 development creating a new business district', 'Minutes from King\'s Cross and Bloomsbury'],
        ['Euston', 'Warren Street', 'Euston Square', 'Mornington Crescent'],
    ),
    'soho': (
        'Soho is London\'s most vibrant creative quarter — home to advertising agencies, film production companies, music labels, and cutting-edge digital businesses. Its central location and buzzing atmosphere make it the address of choice for creative industries.',
        ['London\'s creative and entertainment epicentre', 'Multiple Tube lines: Northern, Central, Bakerloo, Jubilee', 'Unique mix of independent businesses, studios, and agencies', 'Exceptional restaurant and nightlife scene'],
        ['Tottenham Court Road', 'Piccadilly Circus', 'Oxford Circus', 'Leicester Square'],
    ),
    'holborn': (
        'Holborn is a prime legal and professional services location, sitting at the crossroads of the City and the West End. It is home to London\'s four Inns of Court and many of the world\'s top law firms.',
        ['Home to London\'s legal community and major law firms', 'Central line and Piccadilly line connections', 'Between the City and the West End — walk to both', 'Excellent choice of serviced offices and law-focused spaces'],
        ['Holborn', 'Chancery Lane', 'Temple', 'Farringdon'],
    ),
    'north london': (
        'North London offers a compelling mix of vibrant residential areas, strong transport links, and growing business districts. From Islington\'s tech scene to Camden\'s creative energy, north London suits ambitious businesses seeking an alternative to the centre.',
        ['Victoria, Piccadilly, and Northern lines all serve north London', 'Growing tech and media presence in Islington and Camden', 'More affordable than central London with strong talent pools', 'Great access to King\'s Cross and the Elizabeth line'],
        ['Angel', 'King\'s Cross', 'Camden Town', 'Highbury & Islington'],
    ),
    'west london': (
        'West London offers prestigious addresses, excellent transport links, and a diverse business community. From Hammersmith\'s media hub to Chiswick\'s business parks, the area blends professional credibility with quality of life.',
        ['Piccadilly and District lines provide fast central access', 'Close to Heathrow for international travel', 'Strong media, tech, and professional services cluster', 'Green spaces and high quality of life for staff'],
        ['Hammersmith', 'Chiswick', 'Ealing Broadway', 'Acton'],
    ),
    'east london': (
        'East London has transformed into one of the most exciting business destinations in the UK. From Shoreditch\'s tech startups to Canary Wharf\'s finance giants, east London offers an unmatched diversity of workspace and talent.',
        ['Elizabeth line dramatically cuts journey times across London', 'Diverse business ecosystem from tech to finance', 'Competitive pricing with modern, design-led spaces', 'Thriving cultural and social scene'],
        ['Stratford', 'Bethnal Green', 'Whitechapel', 'Canary Wharf'],
    ),
    'south london': (
        'South London is increasingly attractive to businesses priced out of the north or west, offering excellent value, improving transport links, and a growing creative and digital scene centred around areas like Brixton and London Bridge.',
        ['Improving transport links via Thameslink and Overground', 'Vibrant areas: Brixton, Peckham, and the South Bank', 'Significantly lower rents than comparable north London locations', 'Strong community of independent and growing businesses'],
        ['Waterloo', 'London Bridge', 'Brixton', 'Clapham Junction'],
    ),
    'london': (
        'London is the world\'s most important business city — a global hub for finance, technology, law, media, and professional services. With over 500 flexible workspace options across dozens of boroughs, myHQ helps you find the perfect fit wherever in London you need to be.',
        ['World-class transport network: 11 Tube lines, Overground, Elizabeth line', 'Home to 40% of Europe\'s top 500 companies', 'Unrivalled access to talent from world-leading universities', 'Every major business district within 30 minutes by Tube'],
        ['Oxford Circus', 'Liverpool Street', 'King\'s Cross', 'London Bridge', 'Canary Wharf'],
    ),
}

def _get_location_data(loc_str):
    """Return (about_text, highlights, stations) for a given location string."""
    loc_lower = (loc_str or '').lower().strip()

    # 1. Exact match
    if loc_lower in _LOCATION_DATA:
        return _LOCATION_DATA[loc_lower]

    # 2. loc is a substring of a key (e.g. "old street" matches "old street")
    #    Prefer longer key matches (more specific) over shorter ones
    matches = [(key, data) for key, data in _LOCATION_DATA.items() if loc_lower in key]
    if matches:
        matches.sort(key=lambda x: -len(x[0]))  # longest key = most specific
        # But prefer exact matches / close matches over generic "london"
        non_generic = [(k, d) for k, d in matches if k != 'london']
        return (non_generic or matches)[0][1]

    # 3. Key is a substring of loc (e.g. "central london" contains "london")
    #    Again prefer specificity
    matches2 = [(key, data) for key, data in _LOCATION_DATA.items()
                if key in loc_lower and key != 'london']
    if matches2:
        matches2.sort(key=lambda x: -len(x[0]))
        return matches2[0][1]

    # 4. Word overlap
    loc_words = set(loc_lower.split())
    best_score, best_data = 0, None
    for key, data in _LOCATION_DATA.items():
        score = sum(1 for w in key.split() if w in loc_words)
        if score > best_score:
            best_score, best_data = score, data
    if best_data:
        return best_data
    # Generic fallback
    return (
        f'{loc_str} is a well-connected London business district offering a strong mix of flexible workspace options. '
        'The area benefits from excellent transport links and a growing professional community.',
        [
            'Good Tube and bus connectivity throughout',
            'Variety of flexible workspace options available',
            'Close to key London business districts',
            'Growing professional services and tech community',
        ],
        [],   # stations will be derived from actual centres
    )


def _about_location_slide(proposal, all_centres, font, logo, map_img_path=None):
    loc_raw = proposal.get('client_location') or ''
    loc = loc_raw if loc_raw and loc_raw.lower() not in ('none', '') else 'London'
    sl = [R(0, 0, W, H, WHITE)]

    # Split layout: left info panel, right = map
    left_w  = 5200000
    map_x   = left_w + 80000
    map_w   = W - map_x
    map_y   = 0
    map_h   = H

    # Map panel — full height right side
    _map_src = map_img_path if (map_img_path and os.path.isfile(map_img_path)) else (MAP_IMG if os.path.isfile(MAP_IMG) else None)
    sl.append(I(_map_src, map_x, map_y, map_w, map_h))
    sl.append(R(map_x, 0, 6000, H, BLUE))   # thin left border on map

    about_text, highlights, stations = _get_location_data(loc)

    # If no stations found, derive them from the shortlisted centres' transport info
    if not stations:
        seen = set()
        for c in all_centres:
            sn = _extract_station_name(c.get('transport', '') or '')
            if sn and sn not in seen:
                stations.append(sn)
                seen.add(sn)

    # Logo — pick white vs colored based on map image brightness at logo position
    sl.append(_small_logo(logo, bg_image=_map_src))

    # Left: title
    sl.append(T(f'About {loc}', M, 200000, left_w - M, 480000,
                22, BLUE, font=font, bold=True))

    sl.append(T(about_text, M, 730000, left_w - M, 900000,
                9, DARK, font=font, wrap=True))

    sl.append(T('KEY HIGHLIGHTS', M, 1720000, left_w - M, 250000,
                8, BLUE, font=font, bold=True))

    for hi, hl in enumerate(highlights[:4]):
        sl.append(T(f'•  {hl}', M + 80000, 1980000 + hi * 290000,
                    left_w - M - 80000, 260000, 9, DARK, font=font))

    # Popular stations section
    sl.append(T('POPULAR STATIONS', M, 3130000, left_w - M, 250000,
                8, BLUE, font=font, bold=True))

    station_x = M
    station_y = 3400000
    row_gap = 290000
    for si, stn in enumerate(stations[:5]):
        sy = station_y + si * row_gap
        # stations list entries: plain name string OR "Name (walk time)"
        import re as _re
        walk_m = _re.search(r'\((\d+\s*min[^)]*)\)', stn)
        walk_label = walk_m.group(1) if walk_m else ''
        stn_clean = _re.sub(r'\s*\([^)]*\)', '', stn).strip()
        badge_path = station_badge_path(stn_clean)
        if badge_path:
            bw, bh = 580000, 174000
            sl.append(I(badge_path, station_x, sy, bw, bh))
            name_x = station_x + bw + 60000
            sl.append(T(stn_clean, name_x, sy - 10000,
                        left_w - M - bw - 60000, 190000, 9, DARK, font=font, bold=True))
            if walk_label:
                sl.append(T(walk_label, name_x, sy + 170000,
                            left_w - M - bw - 60000, 150000, 7.5, GREY, font=font))
        else:
            sl.append(R(station_x, sy + 55000, 12000, 120000, BLUE))
            sl.append(T(stn_clean, station_x + 80000, sy - 10000,
                        left_w - M - 80000, 190000, 9, DARK, font=font, bold=True))
            if walk_label:
                sl.append(T(walk_label, station_x + 80000, sy + 170000,
                            left_w - M - 80000, 150000, 7.5, GREY, font=font))

    # Shortlisted spaces — compact cards with subtle background
    list_y = station_y + len(stations[:5]) * row_gap + 120000
    card_w = left_w - M - 40000
    card_h = 290000
    card_gap = 12000
    badge_w, badge_h = 600000, 180000   # 320×96 → 3.33:1 exact match, no crop
    if list_y < H - 900000:
        sl.append(T('SHORTLISTED SPACES', M, list_y, card_w, 210000,
                    7.5, BLUE, font=font, bold=True))
        row_y = list_y + 230000
        for ni, c in enumerate(all_centres):
            if row_y + card_h > H - 80000:
                break
            name = c.get('name', '—')
            transport_raw = c.get('transport', '') or ''
            walk_t = _extract_walk_time(transport_raw)
            station_n = _extract_station_name(transport_raw)
            # Card background
            sl.append(R(M, row_y, card_w, card_h, '#F8F9FF'))
            # Index number pill
            sl.append(R(M, row_y, 160000, card_h, BLUE))
            sl.append(T(str(ni + 1), M + 30000, row_y + 80000, 100000, 100000,
                        8, WHITE, font=font, bold=True))
            # Space name — leave room for badge on right
            name_w = card_w - 200000 - (badge_w + 60000) - 20000
            sl.append(T(name, M + 185000, row_y + 40000,
                        name_w, 110000, 8, DARK, font=font, bold=True))
            # Walk time under name
            if walk_t:
                sl.append(T(walk_t, M + 185000, row_y + 160000,
                            name_w, 80000, 6.5, GREY, font=font))
            # Tube badge right-aligned inside card, vertically centred
            # raw=True keeps white background so text on the badge stays visible (card is near-white)
            badge = tube_badge_path(transport_raw, raw=True)
            badge_y = row_y + (card_h - badge_h) // 2
            if badge:
                sl.append(I(badge, M + card_w - badge_w - 30000, badge_y, badge_w, badge_h))
            elif station_n:
                sl.append(T(station_n, M + card_w - 540000, row_y + 90000,
                            500000, 120000, 6.5, GREY, font=font))
            row_y += card_h + card_gap

    return sl


def _make_feature_icon(idx, accent_hex):
    """Draw a clean circular feature icon using PIL and return the file path (cached)."""
    key = f'icon_{idx}_{accent_hex.lstrip("#")}'
    cache = os.path.join(BASE_DIR, 'uploads', f'.{key}.png')
    if os.path.exists(cache) and os.path.getsize(cache) > 200:
        return cache
    try:
        from PIL import Image as _PI, ImageDraw as _ID
        size = 120
        img = _PI.new('RGBA', (size, size), (0, 0, 0, 0))
        d = _ID.Draw(img)
        h = accent_hex.lstrip('#')
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        # Filled circle background
        d.ellipse([0, 0, size - 1, size - 1], fill=(r, g, b, 255))
        wc = (255, 255, 255, 255)
        cx, cy = size // 2, size // 2

        if idx == 0:
            # Checkmark
            pts = [(30, 62), (50, 82), (90, 38)]
            d.line(pts, fill=wc, width=10, joint='curve')
        elif idx == 1:
            # Location pin
            r2 = 20
            d.ellipse([cx - r2, cy - 40, cx + r2, cy], fill=wc)
            d.ellipse([cx - 10, cy - 30, cx + 10, cy - 10], fill=(r, g, b, 255))
            d.polygon([(cx - 18, cy - 8), (cx + 18, cy - 8), (cx, cy + 28)], fill=wc)
        elif idx == 2:
            # Zero cost — circle with cross (no-fee symbol)
            d.ellipse([28, 28, 92, 92], outline=wc, width=7)
            # Diagonal slash
            d.line([(42, 42), (78, 78)], fill=wc, width=7)
        elif idx == 3:
            # Handshake (two arcs)
            d.arc([25, 30, 60, 70], start=180, end=0, fill=wc, width=8)
            d.arc([60, 30, 95, 70], start=180, end=0, fill=wc, width=8)
            d.line([(25, 50), (95, 50)], fill=wc, width=3)
        elif idx == 4:
            # People (two circles + bodies)
            d.ellipse([20, 20, 44, 44], fill=wc)
            d.ellipse([76, 20, 100, 44], fill=wc)
            d.ellipse([45, 15, 75, 45], fill=wc)
            d.arc([5, 55, 55, 90], start=200, end=340, fill=wc, width=8)
            d.arc([65, 55, 115, 90], start=200, end=340, fill=wc, width=8)
            d.arc([32, 50, 88, 90], start=200, end=340, fill=wc, width=9)

        img.save(cache, 'PNG')
        return cache
    except Exception:
        return None


def _how_myhq_helps_slide(font, logo, accent_color=BLUE, icon_color=BLUE):
    """How myHQ helps – 3rd from last slide with icons and cropped workspace photo."""
    sl = [R(0, 0, W, H, WHITE)]

    # Title area
    sl.append(T('How myHQ helps?', M, 280000, 5400000, 520000,
                26, BLUE, font=font, bold=True))
    sl.append(T("Here's what makes us the most preferred brand for workspace solutions",
                M, 850000, 5000000, 380000,
                10, GREY, font=font))

    # Divider under title
    sl.append(L(M, 1280000, 5600000, 1280000, color='#E5E7EB', width=1))

    features = [
        ('Hassle-free process',
         'End-to-end support from search to move-in — we handle everything so you can focus on your business.'),
        ('Largest coverage',
         '500+ spaces across London, all vetted, toured, and ready for you to move in immediately.'),
        ('Zero brokerage fee',
         'No hidden charges, no brokerage. You pay only for the space — full transparency, always.'),
        ('Transparent & ethical',
         'Honest pricing, clear contracts, and a dedicated consultant who works in your best interest.'),
        ('For all team sizes',
         'Whether you are a solo founder or a team of 500, we match you to the perfect space.'),
    ]

    left_panel_w = 5600000
    feature_y_start = 1420000
    feature_row_h   = 1000000
    icon_size = 500000   # 500k EMU square icon

    for i, (title, desc) in enumerate(features):
        fy = feature_y_start + i * feature_row_h
        icon_path = _make_feature_icon(i, accent_color)
        if icon_path and os.path.isfile(icon_path):
            sl.append(I(icon_path, M, fy + 30000, icon_size, icon_size))
        else:
            sl.append(R(M, fy + 30000, icon_size, icon_size, accent_color))

        sl.append(T(title, M + icon_size + 80000, fy + 40000,
                    left_panel_w - M - icon_size - 80000, 260000,
                    11, DARK, font=font, bold=True))
        sl.append(T(desc, M + icon_size + 80000, fy + 300000,
                    left_panel_w - M - icon_size - 80000, 460000,
                    9, GREY, font=font, wrap=True))

        if i < len(features) - 1:
            sl.append(L(M + icon_size + 80000, fy + feature_row_h - 60000,
                        left_panel_w, fy + feature_row_h - 60000,
                        color='#F0F0F0', width=1))

    # Right photo panel – cropped workspace image
    photo_x = left_panel_w + 200000
    photo_w = W - photo_x - 100000
    crop_photo = CROP_IMG if os.path.isfile(CROP_IMG) else None
    sl.append(I(crop_photo, photo_x, 0, photo_w, H))

    # Logo appended last so it renders above the photo panel
    sl.append(_small_logo(logo))

    return sl


# ── London slide deck ───────────────────────────────────────────────────────────

def build_london_slides(proposal, db_centres, manual_centres):
    F       = 'Mont'
    logo    = get_logo_png()
    logo_w  = get_logo_png(white=True)
    all_centres = db_centres + manual_centres
    slides = []

    # ── Slide 1: Cover ──────────────────────────────────────────────────────────
    sl = []
    # Left panel (blush)
    sl.append(R(0, 0, 5100000, H, BLUSH))
    # Right photo (city skyline, high-res)
    cover_photo = COVER_IMG if os.path.isfile(COVER_IMG) else None
    sl.append(I(cover_photo, 5100000, 0, W - 5100000, H))
    # myHQ logo top-left composited on blush panel
    sl.append(I(get_logo_png(white=False), M, 280000, 1700000, 637000))
    # Title
    client = proposal.get('client_company') or proposal.get('client_name') or 'Your Company'
    sl.append(T(f'Workspace Proposal\nfor {client}',
                M, 2100000, 4600000, 1800000,
                30, BLUE, font=F, bold=True))
    # Subtitle
    sl.append(T('Flexible workspaces, transparent pricing, no brokerage',
                M, 3950000, 4600000, 450000,
                11, GREY, font=F))
    # Client name bottom
    sl.append(T(client, M, 5650000, 4600000, 450000,
                16, BLUE, font=F, bold=True))
    slides.append(sl)

    # ── Slide 2: Comparison Table ────────────────────────────────────────────────
    slides.append(
        _comparison_table_slide(all_centres,
                                header_fill=AQUA,
                                header_text_color=DARK,
                                font=F, logo=logo))

    # ── Slide 3: Client Requirements ────────────────────────────────────────────
    slides.append(
        _client_requirements_slide(
            proposal, 'london', db_centres, manual_centres,
            font=F, logo=logo,
            label_fill=AQUA,
            label_text=DARK,
            date_label_color=MINT))

    # ── Slide 4: About Location ──────────────────────────────────────────────────
    _map_tmp = os.path.join(BASE_DIR, 'uploads', 'proposals', f'map_{id(proposal)}.png')
    _map_path = generate_proposal_map(all_centres, _map_tmp) or MAP_IMG
    slides.append(_about_location_slide(proposal, all_centres, F, logo, map_img_path=_map_path))

    # ── Slides 5+: Centre slides ─────────────────────────────────────────────────
    for idx, centre in enumerate(all_centres, 1):
        slides.append(_centre_slide(idx, centre, 'london', F, logo))

    # ── 3rd from last: How myHQ helps ───────────────────────────────────────────
    slides.append(_how_myhq_helps_slide(F, logo, accent_color=AQUA, icon_color=BLUE))

    # ── Existing Clients ────────────────────────────────────────────────────────
    slides.append(_existing_clients_slide(F, logo))

    # ── Testimonials ────────────────────────────────────────────────────────────
    slides.append(_testimonials_slide(F, logo, bg_color=WHITE, heading_color=DARK))

    return slides


# ── India slide deck ────────────────────────────────────────────────────────────

def build_india_slides(proposal, db_centres, manual_centres):
    F       = 'Mont'
    logo    = get_logo_png()
    logo_w  = get_logo_png(white=True)
    all_centres = db_centres + manual_centres
    slides = []

    # ── Slide 1: Cover ──────────────────────────────────────────────────────────
    sl = []
    sl.append(R(0, 0, 5100000, H, WHITE))
    cover_photo = COVER_IMG if os.path.isfile(COVER_IMG) else None
    sl.append(I(cover_photo, 5100000, 0, W - 5100000, H))
    sl.append(I(get_logo_png(white=False), M, 280000, 1700000, 637000))
    sl.append(T("The search for your\nperfect workspace\nends here",
                M, 2100000, 4600000, 1800000,
                30, BLUE, font=F, bold=True))
    sl.append(T("India's leading marketplace for flexible workspace solutions",
                M, 3950000, 4600000, 450000,
                11, GREY, font=F))
    client = proposal.get('client_company') or proposal.get('client_name') or 'Your Company'
    sl.append(T(client, M, 5650000, 4600000, 450000,
                16, DARK, font=F, bold=True))
    slides.append(sl)

    # ── Slide 2: Comparison Table ────────────────────────────────────────────────
    slides.append(
        _comparison_table_slide(all_centres,
                                header_fill=BLUE,
                                header_text_color=WHITE,
                                font=F, logo=logo))

    # ── Slide 3: Client Requirements ────────────────────────────────────────────
    slides.append(
        _client_requirements_slide(
            proposal, 'india', db_centres, manual_centres,
            font=F, logo=logo,
            label_fill=BLUE,
            label_text=WHITE,
            date_label_color='#AAAAFF'))

    # ── Slide 4: About Location ──────────────────────────────────────────────────
    _map_tmp = os.path.join(BASE_DIR, 'uploads', 'proposals', f'map_{id(proposal)}.png')
    _map_path = generate_proposal_map(all_centres, _map_tmp) or MAP_IMG
    slides.append(_about_location_slide(proposal, all_centres, F, logo, map_img_path=_map_path))

    # ── Slides 5+: Centre slides ─────────────────────────────────────────────────
    for idx, centre in enumerate(all_centres, 1):
        slides.append(_centre_slide(idx, centre, 'india', F, logo))

    # ── 3rd from last: How myHQ helps ───────────────────────────────────────────
    slides.append(_how_myhq_helps_slide(F, logo, accent_color=BLUE, icon_color=BLUE))

    # ── Existing Clients ────────────────────────────────────────────────────────
    slides.append(_existing_clients_slide(F, logo))

    # ── Testimonials ────────────────────────────────────────────────────────────
    slides.append(_testimonials_slide(F, logo, bg_color=WHITE, heading_color=DARK))

    return slides


# ── Bold/Creative (3rd) slide deck ─────────────────────────────────────────────

def _bold_centre_slide(idx, centre, font, logo):
    """Centre detail slide — Trampery-style: dark navy left panel + full-height photo right."""
    NAVY  = '#0B0E2B'
    PANEL = '#111740'

    sl = [R(0, 0, W, H, NAVY)]

    images = centre.get('images', [])
    logo_w = get_logo_png(white=True)

    # ── Right: full-height hero photo ────────────────────────────────────────
    panel_w = int(W * 0.43)
    photo_x = panel_w
    photo_w = W - panel_w
    hero = images[0] if images else None
    sl.append(R(photo_x, 0, photo_w, H, '#1A1E40'))
    if hero:
        sl.append(I(hero, photo_x, 0, photo_w, H))

    # AQUA vertical divider
    sl.append(R(panel_w - 24000, 0, 24000, H, AQUA))

    # Logo — dynamic: white on dark hero, colored on light hero
    sl.append(_small_logo(logo_w, bg_image=hero))

    # ── Secondary image strip at bottom of photo ─────────────────────────────
    if len(images) > 1:
        strip_h = int(H * 0.17)
        strip_y = H - strip_h
        n_thumbs = min(len(images) - 1, 3)
        thumb_w  = (photo_w - (n_thumbs - 1) * 20000) // n_thumbs
        for ti, img in enumerate(images[1:n_thumbs + 1]):
            tx = photo_x + ti * (thumb_w + 20000)
            sl.append(I(img, tx, strip_y, thumb_w, strip_h))
        sl.append(R(photo_x, strip_y - 10000, photo_w, 10000, AQUA))

    # ── Left dark panel ───────────────────────────────────────────────────────
    # Large decorative index number (ghost, very dark)
    sl.append(T(f'{idx:02d}', M, 80000, 800000, 750000, 68, '#1E2450', font=font, bold=True))

    # Centre name
    name = centre.get('name', 'Centre')
    brand = (centre.get('brand') or '').strip()
    display_name = f'{brand} – {name}' if brand and brand.lower() not in name.lower() else name

    sl.append(T(display_name.upper(),
                M, 870000, panel_w - M - 120000, 700000,
                18, WHITE, font=font, bold=True, wrap=True))

    # AQUA rule under name
    sl.append(R(M, 1650000, 700000, 10000, AQUA))

    cy = 1720000
    info_w = panel_w - M - 120000

    # Address
    sl.append(T('ADDRESS', M, cy, info_w, 170000, 7, '#6B80A0', font=font, bold=True))
    cy += 180000
    sl.append(T(centre.get('address', '—'), M, cy, info_w, 280000, 9, '#BCCDE0', font=font, wrap=True))
    cy += 330000

    # Price + Hours row
    half = (info_w - 80000) // 2
    sl.append(T('PRICE / SEAT', M, cy, half, 170000, 7, '#6B80A0', font=font, bold=True))
    sl.append(T('OPEN HOURS', M + half + 80000, cy, half, 170000, 7, '#6B80A0', font=font, bold=True))
    cy += 180000
    sl.append(T(price_str(centre), M, cy, half, 260000, 13, AQUA, font=font, bold=True))
    sl.append(T(centre.get('open_hours') or '9:00 AM – 6:00 PM',
                M + half + 80000, cy, half, 260000, 9, WHITE, font=font))
    cy += 320000

    # Nearest tube
    transport_raw = centre.get('transport', '') or ''
    station_n = _extract_station_name(transport_raw)
    walk_t    = _extract_walk_time(transport_raw)
    badge     = tube_badge_path(transport_raw)

    sl.append(T('NEAREST TUBE', M, cy, info_w, 170000, 7, '#6B80A0', font=font, bold=True))
    cy += 180000
    if badge:
        # Pale navy pill behind the badge so it sits cleanly on the dark panel
        sl.append(R(M - 30000, cy - 20000, 640000, 214000, '#1C2448'))
        sl.append(I(badge, M, cy, 580000, 174000))
        sl.append(T(station_n or '', M + 640000, cy, info_w - 640000, 180000, 9, WHITE, font=font, bold=True))
        if walk_t:
            sl.append(T(walk_t, M + 640000, cy + 185000, info_w - 640000, 160000, 7.5, '#6B80A0', font=font))
        cy += 360000
    elif station_n:
        sl.append(R(M, cy + 40000, 12000, 130000, AQUA))
        sl.append(T(station_n, M + 60000, cy, info_w - 60000, 190000, 9, WHITE, font=font, bold=True))
        if walk_t:
            sl.append(T(walk_t, M + 60000, cy + 190000, info_w - 60000, 160000, 7.5, '#6B80A0', font=font))
        cy += 360000

    cy += 30000
    sl.append(R(M, cy, 600000, 6000, '#2A3A5A'))
    cy += 50000

    # Amenities — 2-column on dark background
    sl.append(T('AMENITIES', M, cy, info_w, 170000, 7, '#6B80A0', font=font, bold=True))
    cy += 185000

    try:
        ams_raw = centre.get('amenities', '[]')
        ams = json.loads(ams_raw) if isinstance(ams_raw, str) else (ams_raw or [])
    except Exception:
        ams = []

    cols   = 2
    col_w2 = info_w // cols
    row_h2 = 270000
    dot_sz = 90000
    dot_gp = 65000
    txt_h2 = 190000

    for ai, slug in enumerate(ams[:8]):
        arow = ai // cols
        if cy + arow * row_h2 + row_h2 > H - 80000:
            break
        slug = str(slug).lower().strip()
        label = AMENITY_LABELS.get(slug, slug.replace('_', ' ').replace('-', ' ').title())
        if len(label) > 20:
            label = label[:19] + '…'
        acol = ai % cols
        ax   = M + acol * col_w2
        ay   = cy + arow * row_h2
        dot_y = ay + (row_h2 - dot_sz) // 2
        txt_y = ay + (row_h2 - txt_h2) // 2
        color = _AMENITY_COLORS[ai % len(_AMENITY_COLORS)]
        sl.append(R(ax, dot_y, dot_sz, dot_sz, color))
        sl.append(T(label, ax + dot_sz + dot_gp, txt_y,
                    col_w2 - dot_sz - dot_gp - 40000, txt_h2,
                    7.5, '#BCCDE0', font=font, wrap=False))

    # AQUA bottom accent bar
    sl.append(R(0, H - 55000, panel_w, 55000, AQUA))

    return sl


def _bold_about_location_slide(proposal, all_centres, font, logo, map_img_path=None):
    """About Location slide — Trampery-style dark navy + AQUA accents."""
    NAVY = '#0B0E2B'
    loc  = proposal.get('client_location') or 'London'
    about_text, highlights, stations = _get_location_data(loc)

    sl = [R(0, 0, W, H, NAVY)]

    left_w = W // 2 - 100000

    # Map — full height right half, no logo on top of it
    map_x = left_w + 60000
    map_w = W - map_x
    map_y = 0
    map_h = H
    _map_src = map_img_path if (map_img_path and os.path.isfile(map_img_path)) else \
               (MAP_IMG if os.path.isfile(MAP_IMG) else None)
    sl.append(I(_map_src, map_x, map_y, map_w, map_h))
    sl.append(R(map_x, 0, 6000, H, AQUA))   # thin AQUA divider

    # AQUA accent bar at top of left panel only
    sl.append(R(0, 0, left_w, 80000, AQUA))

    # Title
    sl.append(T(f'About {loc}'.upper(), M, 160000, left_w - M, 420000,
                22, WHITE, font=font, bold=True))

    sl.append(T(about_text, M, 680000, left_w - M, 800000,
                9, '#BCCDE0', font=font, wrap=True))

    sl.append(T('KEY HIGHLIGHTS', M, 1580000, left_w - M, 220000,
                7.5, AQUA, font=font, bold=True))
    sl.append(R(M, 1820000, 700000, 8000, AQUA))

    for hi, hl in enumerate(highlights[:4]):
        sl.append(R(M, 1910000 + hi * 290000 + 90000, 10000, 110000, AQUA))
        sl.append(T(hl, M + 80000, 1910000 + hi * 290000,
                    left_w - M - 80000, 260000, 9, '#BCCDE0', font=font))

    sl.append(T('POPULAR STATIONS', M, 3110000, left_w - M, 220000,
                7.5, AQUA, font=font, bold=True))
    sl.append(R(M, 3340000, 700000, 8000, AQUA))

    import re as _re
    station_y = 3410000
    row_gap   = 290000
    for si, stn in enumerate(stations[:5]):
        sy = station_y + si * row_gap
        walk_m = _re.search(r'\((\d+\s*min[^)]*)\)', stn)
        walk_label = walk_m.group(1) if walk_m else ''
        stn_clean  = _re.sub(r'\s*\([^)]*\)', '', stn).strip()
        badge_path = station_badge_path(stn_clean)
        if badge_path:
            sl.append(I(badge_path, M, sy, 580000, 174000))
            sl.append(T(stn_clean, M + 620000, sy - 8000,
                        left_w - M - 620000, 190000, 9, WHITE, font=font, bold=True))
            if walk_label:
                sl.append(T(walk_label, M + 620000, sy + 165000,
                            left_w - M - 620000, 150000, 7.5, '#6B80A0', font=font))
        else:
            sl.append(R(M, sy + 55000, 12000, 120000, AQUA))
            sl.append(T(stn_clean, M + 80000, sy - 8000,
                        left_w - M - 80000, 190000, 9, WHITE, font=font, bold=True))

    # Shortlisted spaces
    list_y = station_y + len(stations[:5]) * row_gap + 150000
    if list_y < H - 1000000:
        sl.append(T('SHORTLISTED SPACES', M, list_y, left_w - M, 220000,
                    7.5, AQUA, font=font, bold=True))
        row_y = list_y + 260000
        for ni, c in enumerate(all_centres):
            if row_y > H - 320000:
                break
            name = c.get('name', '—')
            sl.append(T(f'{ni + 1}.  {name}', M + 80000, row_y,
                        left_w - M - 1100000, 250000, 9, '#BCCDE0', font=font))
            row_y += 270000

    # Logo — top-right corner, dark navy pill so it sits on map cleanly
    lw = get_logo_png(white=True)
    sl.append(R(W - 1610000, 80000, 1650000, 560000, NAVY))
    sl.append(I(lw, W - 1550000, 130000, 1100000, 412000))
    return sl


def build_bold_slides(proposal, db_centres, manual_centres):
    """3rd creative template — Trampery-inspired: dark navy panels, full-bleed photos, AQUA accents."""
    F       = 'Mont'
    logo    = get_logo_png()
    logo_w  = get_logo_png(white=True)
    all_centres = db_centres + manual_centres
    slides  = []
    NAVY    = '#0B0E2B'

    # ── Cover slide ─────────────────────────────────────────────────────────────
    sl = []
    cover_photo = COVER_IMG if os.path.isfile(COVER_IMG) else None
    # Full-bleed cover photo right side
    sl.append(R(0, 0, W, H, NAVY))
    sl.append(I(cover_photo, int(W * 0.44), 0, int(W * 0.56), H))
    # Dark panel left
    sl.append(R(0, 0, int(W * 0.46), H, NAVY))
    # AQUA vertical divider
    sl.append(R(int(W * 0.44) - 24000, 0, 24000, H, AQUA))
    # AQUA top stripe full width
    sl.append(R(0, 0, W, 75000, AQUA))
    # White logo top-left
    sl.append(I(logo_w, M, 260000, 1600000, 599000))
    # Huge title
    client = proposal.get('client_company') or proposal.get('client_name') or 'Your Company'
    sl.append(T('WORKSPACE\nPROPOSAL',
                M, 1600000, int(W * 0.44) - M - 100000, 2400000,
                50, WHITE, font=F, bold=True))
    sl.append(T(f'PREPARED FOR  {client.upper()}',
                M, 4200000, int(W * 0.44) - M - 100000, 500000,
                10, AQUA, font=F, bold=True))
    sl.append(T('Flexible workspaces  ·  Zero brokerage  ·  Expert guidance',
                M, 4790000, int(W * 0.44) - M - 100000, 380000,
                8.5, '#4A5A7A', font=F))
    # AQUA bottom accent bar
    sl.append(R(0, H - 60000, int(W * 0.46), 60000, AQUA))
    slides.append(sl)

    # ── Comparison table (BLUE header) ──────────────────────────────────────────
    slides.append(
        _comparison_table_slide(all_centres,
                                header_fill=BLUE,
                                header_text_color=WHITE,
                                font=F, logo=logo))

    # ── Client requirements (AQUA labels) ───────────────────────────────────────
    slides.append(
        _client_requirements_slide(
            proposal, 'london', db_centres, manual_centres,
            font=F, logo=logo,
            label_fill=AQUA,
            label_text=DARK,
            date_label_color=MINT))

    # ── About location ───────────────────────────────────────────────────────────
    _map_tmp = os.path.join(BASE_DIR, 'uploads', 'proposals', f'map_bold_{id(proposal)}.png')
    _map_path = generate_proposal_map(all_centres, _map_tmp) or MAP_IMG
    slides.append(_bold_about_location_slide(proposal, all_centres, F, logo, map_img_path=_map_path))

    # ── Centre slides (dark panel layout) ───────────────────────────────────────
    for idx, centre in enumerate(all_centres, 1):
        slides.append(_bold_centre_slide(idx, centre, F, logo))

    # ── How myHQ helps ───────────────────────────────────────────────────────────
    slides.append(_how_myhq_helps_slide(F, logo, accent_color=AQUA, icon_color=BLUE))

    # ── Existing clients ─────────────────────────────────────────────────────────
    slides.append(_existing_clients_slide(F, logo))

    # ── Testimonials ─────────────────────────────────────────────────────────────
    slides.append(_testimonials_slide(F, logo, bg_color=WHITE, heading_color=DARK))

    return slides


# ── Template auto-detection ─────────────────────────────────────────────────────

def detect_template_features(pptx_path):
    """Analyse a PPTX template and return a dict describing its features."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
        prs = Presentation(pptx_path)
        slides_info = []
        for sli, slide in enumerate(prs.slides):
            img_placeholders = 0
            text_placeholders = 0
            has_title = False
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_placeholders += 1
                elif hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        has_title = True
                    elif ph_type == PP_PLACEHOLDER.PICTURE:
                        img_placeholders += 1
                    else:
                        text_placeholders += 1
                elif hasattr(shape, 'text') and shape.text.strip():
                    text_placeholders += 1
            slides_info.append({
                'slide': sli + 1,
                'images': img_placeholders,
                'texts': text_placeholders,
                'has_title': has_title,
            })
        # Heuristic: the slide with the most image placeholders is the centre slide
        max_imgs = max((s['images'] for s in slides_info), default=0)
        # Slides with max images are likely centre slides
        centre_slides = [s for s in slides_info if s['images'] == max_imgs and max_imgs > 0]
        images_per_centre = max_imgs if centre_slides else 4
        return {
            'images_per_centre': images_per_centre,
            'total_slides': len(prs.slides),
            'slides_detail': slides_info,
        }
    except Exception as e:
        return {'images_per_centre': 4, 'total_slides': 0, 'error': str(e)}


# ── PPTX renderer ──────────────────────────────────────────────────────────────

def _rgb(hex_str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def render_pptx(slides, out_path):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu as EMU

    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    blank_layout = prs.slide_layouts[6]  # completely blank

    ALIGN_MAP = {
        'L': PP_ALIGN.LEFT,
        'C': PP_ALIGN.CENTER,
        'R': PP_ALIGN.RIGHT,
    }

    for slide_cmds in slides:
        sl = prs.slides.add_slide(blank_layout)
        for cmd in slide_cmds:
            ct = cmd['cmd']

            if ct == 'rect':
                x = Emu(cmd['x']); y = Emu(cmd['y'])
                w = Emu(cmd['w']); h = Emu(cmd['h'])
                sh = sl.shapes.add_shape(1, x, y, w, h)
                sh.line.fill.background()
                sh.fill.solid()
                sh.fill.fore_color.rgb = _rgb(cmd['fill'])

            elif ct == 'text':
                if not str(cmd.get('text', '')).strip():
                    continue
                x = Emu(cmd['x']); y = Emu(cmd['y'])
                w = Emu(cmd['w']); h = Emu(cmd['h'])
                tb = sl.shapes.add_textbox(x, y, w, h)
                tf = tb.text_frame
                tf.word_wrap = cmd.get('wrap', True)
                align = ALIGN_MAP.get(cmd.get('align', 'L'), PP_ALIGN.LEFT)
                lines = cmd['text'].split('\n')
                first = True
                for line in lines:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.alignment = align
                    run = p.add_run()
                    run.text = line
                    run.font.name    = cmd.get('font', 'Mont')
                    run.font.size    = Pt(cmd['size'])
                    run.font.color.rgb = _rgb(cmd['color'])
                    run.font.bold    = cmd.get('bold', False)
                    run.font.italic  = cmd.get('italic', False)

            elif ct == 'img':
                x = Emu(cmd['x']); y = Emu(cmd['y'])
                w = Emu(cmd['w']); h = Emu(cmd['h'])
                path = cmd.get('path')
                # Download remote URL to local cache if needed
                if path and str(path).startswith('http'):
                    import hashlib as _hlib, urllib.request as _ur2, ssl as _ssl2
                    _cache_dir = os.path.join(BASE_DIR, 'uploads', '.img_cache')
                    os.makedirs(_cache_dir, exist_ok=True)
                    _ext = str(path).split('?')[0].rsplit('.', 1)[-1][:4] or 'jpg'
                    _cache = os.path.join(_cache_dir, _hlib.md5(str(path).encode()).hexdigest() + '.' + _ext)
                    if not os.path.isfile(_cache):
                        try:
                            _ctx2 = _ssl2._create_unverified_context()
                            _req2 = _ur2.Request(str(path), headers={'User-Agent': 'myHQ/1.0'})
                            with _ur2.urlopen(_req2, timeout=12, context=_ctx2) as _r2:
                                with open(_cache, 'wb') as _cf:
                                    _cf.write(_r2.read())
                        except Exception:
                            path = None
                    if path:
                        path = _cache
                if path and os.path.isfile(str(path)):
                    try:
                        from PIL import Image as _PILImg
                        import io as _io2
                        _im = _PILImg.open(str(path))
                        # Normalise palette/odd modes to PNG-safe RGB or RGBA
                        if _im.mode == 'P':
                            _im = _im.convert('RGBA')
                        use_path = str(path)
                        # If image was changed, write to a temp buffer and use that
                        if _im.mode not in ('RGB', 'RGBA', 'L'):
                            _im = _im.convert('RGB')
                        if _im.mode == 'RGBA':
                            # Keep for PPTX (transparency is OK), but re-save to fix palette artefacts
                            _buf = _io2.BytesIO()
                            _im.save(_buf, 'PNG')
                            _buf.seek(0)
                            iw, ih = _im.size
                        else:
                            iw, ih = _im.size
                        tw, th = cmd['w'], cmd['h']
                        # Cover-crop fractions
                        scale = max(tw / iw, th / ih)
                        sw_px = iw * scale
                        sh_px = ih * scale
                        cx = max(0.0, (sw_px - tw) / (2 * sw_px))
                        cy = max(0.0, (sh_px - th) / (2 * sh_px))
                        if _im.mode == 'RGBA':
                            pic = sl.shapes.add_picture(_buf, x, y, w, h)
                        else:
                            pic = sl.shapes.add_picture(use_path, x, y, w, h)
                        pic.crop_left   = cx
                        pic.crop_right  = cx
                        pic.crop_top    = cy
                        pic.crop_bottom = cy
                        continue
                    except Exception:
                        pass
                # grey placeholder
                sh = sl.shapes.add_shape(1, x, y, w, h)
                sh.line.fill.background()
                sh.fill.solid()
                sh.fill.fore_color.rgb = _rgb('#DDDDDD')

            elif ct == 'line':
                x1 = Emu(cmd['x1']); y1 = Emu(cmd['y1'])
                x2 = Emu(cmd['x2']); y2 = Emu(cmd['y2'])
                # Represent as a thin rectangle
                lw = max(Emu(cmd.get('width', 1) * 12700), Emu(12700))
                if x1 == x2:
                    sh = sl.shapes.add_shape(1, x1, y1, lw, Emu(abs(cmd['y2'] - cmd['y1'])))
                else:
                    sh = sl.shapes.add_shape(1, x1, y1, Emu(abs(cmd['x2'] - cmd['x1'])), lw)
                sh.line.fill.background()
                sh.fill.solid()
                sh.fill.fore_color.rgb = _rgb(cmd.get('color', '#E5E7EB'))

    prs.save(out_path)
    return out_path


# ── PDF renderer ────────────────────────────────────────────────────────────────

def render_pdf(slides, out_path):
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.colors import HexColor
    from reportlab.lib.utils import simpleSplit

    PW, PH  = 960.0, 540.0   # points (13.33" × 7.5" @ 72 dpi)
    E2P     = 1.0 / 12700.0  # EMU → pt

    def ep(v):
        return float(v) * E2P

    def hex_color(h):
        return HexColor(f'#{h.lstrip("#")}')

    def wrap_lines(text, fn, fs, max_w):
        lines = []
        for para in text.split('\n'):
            try:
                wrapped = simpleSplit(para, fn, fs, max_w)
                lines.extend(wrapped if wrapped else [''])
            except Exception:
                lines.extend(para.split('\n') or [''])
        return lines

    c = rl_canvas.Canvas(out_path, pagesize=(PW, PH))

    for slide_cmds in slides:
        for cmd in slide_cmds:
            ct = cmd['cmd']
            sx  = ep(cmd.get('x', 0))
            sy  = ep(cmd.get('y', 0))   # top in EMU space (pt)
            sw  = ep(cmd.get('w', 0))
            sh  = ep(cmd.get('h', 0))
            # PDF y=0 is bottom-left; our y is measured from top
            pdf_y = PH - sy - sh        # bottom of box in PDF coords

            if ct == 'rect':
                col = hex_color(cmd['fill'])
                c.setFillColor(col)
                c.setStrokeColor(col)
                c.rect(sx, pdf_y, sw, sh, fill=1, stroke=0)

            elif ct == 'text':
                text = str(cmd.get('text', ''))
                if not text.strip():
                    continue
                fs     = float(cmd['size'])
                bold   = cmd.get('bold', False)
                italic = cmd.get('italic', False)
                if bold and italic:   fn = 'Helvetica-BoldOblique'
                elif bold:            fn = 'Helvetica-Bold'
                elif italic:          fn = 'Helvetica-Oblique'
                else:                 fn = 'Helvetica'

                c.setFont(fn, fs)
                c.setFillColor(hex_color(cmd['color']))
                align  = cmd.get('align', 'L')
                do_wrap = cmd.get('wrap', True)

                if do_wrap:
                    lines = wrap_lines(text, fn, fs, sw - 4)
                else:
                    lines = text.split('\n')

                line_h = fs * 1.35
                # Start at near the top of the box
                text_top = PH - sy - fs * 0.85
                for li, line in enumerate(lines):
                    ty = text_top - li * line_h
                    if ty < pdf_y:          # clipped
                        break
                    if align == 'C':
                        c.drawCentredString(sx + sw / 2, ty, line)
                    elif align == 'R':
                        c.drawRightString(sx + sw, ty, line)
                    else:
                        c.drawString(sx + 2, ty, line)

            elif ct == 'img':
                path = cmd.get('path')
                # Download remote URL to local cache if needed
                if path and str(path).startswith('http'):
                    import hashlib as _hlib, urllib.request as _ur2, ssl as _ssl2
                    _cache_dir = os.path.join(BASE_DIR, 'uploads', '.img_cache')
                    os.makedirs(_cache_dir, exist_ok=True)
                    _ext = str(path).split('?')[0].rsplit('.', 1)[-1][:4] or 'jpg'
                    _cache = os.path.join(_cache_dir, _hlib.md5(str(path).encode()).hexdigest() + '.' + _ext)
                    if not os.path.isfile(_cache):
                        try:
                            _ctx2 = _ssl2._create_unverified_context()
                            _req2 = _ur2.Request(str(path), headers={'User-Agent': 'myHQ/1.0'})
                            with _ur2.urlopen(_req2, timeout=12, context=_ctx2) as _r2:
                                with open(_cache, 'wb') as _cf:
                                    _cf.write(_r2.read())
                        except Exception:
                            path = None
                    if path:
                        path = _cache
                if path and os.path.isfile(str(path)):
                    try:
                        from PIL import Image as _PILImg
                        import io as _io
                        _raw = _PILImg.open(str(path))

                        # Detect meaningful transparency (logos, badges, stripped PNGs)
                        has_alpha = False
                        if _raw.mode in ('RGBA', 'LA', 'P'):
                            _test = _raw.convert('RGBA')
                            alpha = _test.split()[3]
                            has_alpha = min(alpha.getdata()) < 200

                        if has_alpha:
                            # Preserve transparency — draw as PNG with mask='auto'
                            _rgba = _raw.convert('RGBA')
                            iw, ih = _rgba.size
                            target_ratio = sw / sh if sh > 0 else 1.0
                            src_ratio    = iw / ih if ih > 0 else 1.0
                            # Only cover-crop if ratio mismatch > 5% to avoid cutting badges
                            if src_ratio > target_ratio * 1.05:
                                new_w = int(ih * target_ratio)
                                x0    = (iw - new_w) // 2
                                _rgba = _rgba.crop((x0, 0, x0 + new_w, ih))
                            elif src_ratio < target_ratio * 0.95:
                                new_h = int(iw / target_ratio)
                                y0    = (ih - new_h) // 2
                                _rgba = _rgba.crop((0, y0, iw, y0 + new_h))
                            buf = _io.BytesIO()
                            _rgba.save(buf, 'PNG')
                            buf.seek(0)
                            ir = ImageReader(buf)
                            c.drawImage(ir, sx, pdf_y, sw, sh,
                                        preserveAspectRatio=False, mask='auto')
                            continue

                        # Opaque image — composite any RGBA onto white, save as JPEG
                        if _raw.mode in ('P', 'RGBA', 'LA'):
                            _raw = _raw.convert('RGBA')
                            bg   = _PILImg.new('RGB', _raw.size, (255, 255, 255))
                            bg.paste(_raw, mask=_raw.split()[3])
                            _im  = bg
                        else:
                            _im = _raw.convert('RGB')
                        iw, ih = _im.size
                        # Cover-crop to target aspect ratio
                        target_ratio = sw / sh if sh > 0 else 1.0
                        src_ratio    = iw / ih if ih > 0 else 1.0
                        if src_ratio > target_ratio:
                            new_w = int(ih * target_ratio)
                            x0    = (iw - new_w) // 2
                            _im   = _im.crop((x0, 0, x0 + new_w, ih))
                        elif src_ratio < target_ratio:
                            new_h = int(iw / target_ratio)
                            y0    = (ih - new_h) // 2
                            _im   = _im.crop((0, y0, iw, y0 + new_h))
                        buf = _io.BytesIO()
                        _im.save(buf, 'JPEG', quality=90)
                        buf.seek(0)
                        ir = ImageReader(buf)
                        c.drawImage(ir, sx, pdf_y, sw, sh,
                                    preserveAspectRatio=False, mask='auto')
                        continue
                    except Exception:
                        pass
                # grey placeholder
                c.setFillColor(HexColor('#DDDDDD'))
                c.setStrokeColor(HexColor('#CCCCCC'))
                c.rect(sx, pdf_y, sw, sh, fill=1, stroke=1)

            elif ct == 'line':
                x1 = ep(cmd['x1']); y1 = PH - ep(cmd['y1'])
                x2 = ep(cmd['x2']); y2 = PH - ep(cmd['y2'])
                lw = float(cmd.get('width', 1))
                c.setStrokeColor(hex_color(cmd.get('color', '#E5E7EB')))
                c.setLineWidth(lw)
                c.line(x1, y1, x2, y2)

        c.showPage()

    c.save()
    return out_path
